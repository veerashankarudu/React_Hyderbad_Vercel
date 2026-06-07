#!/usr/bin/env python3
"""Generate QuizHub AI Demo PowerPoint with screenshots and voiceover speaker notes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# --- Configuration ---
SHOTS_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(SHOTS_DIR, "QuizHub_AI_Demo.pptx")

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# --- Slide Data: (image_file, title, bullet_points, speaker_notes) ---
SLIDES = [
    # Title Slide
    (None, None, None, None),  # handled specially

    # Scene 1: Login
    ("00-login-page.png",
     "Role-Based Login",
     ["JWT-secured authentication", "SME & Admin roles", "Quick-access demo credentials", "Password reset via email OTP"],
     "Welcome to QuizHub AI — an intelligent MCQ management platform built for Accenture's Hack-N-Stack 2026. This platform enables Subject Matter Experts to create, review, and manage technical interview questions powered by AI. We log in as an SME. The platform supports role-based access with SME and Admin roles, secured with JWT authentication."),

    # Scene 2: Dashboard
    ("01-dashboard.png",
     "SME Dashboard",
     ["Question stats at a glance", "Quick-access action cards", "Recent activity feed", "Tech stack distribution chart"],
     "The SME dashboard shows a summary of your questions — total created, pending reviews, approved count, and recent activity. Quick-access cards let you create MCQs, view your questions, and access AI tools."),

    # Scene 3: Create MCQ
    ("03-create-mcq.png",
     "Create MCQ — Manual Authoring",
     ["Select tech stack & topic", "Set difficulty (Easy/Medium/Hard)", "Rich text question stem with code support", "4 options with correct answer marking"],
     "Creating a question is straightforward. Select a tech stack like Spring Boot, choose a topic, set the difficulty level, write the question stem, and provide four options with the correct answer marked. You can save as draft or submit directly for review. The platform supports rich text, code snippets, and image-based questions."),

    # Scene 4: AI Studio
    ("11-ai-studio.png",
     "AI-Powered Question Generation",
     ["Spring AI + LLM integration", "Generate multiple MCQs at once", "Auto-generated distractors & explanations", "Edit before submission"],
     "One of our most powerful features — AI-powered question generation. Select a tech stack and topic, choose difficulty and count, and the AI generates complete MCQs with distractors, correct answers, and explanations. The AI uses Spring AI integrated with large language models, generating technically accurate questions that test real understanding."),

    # Scene 5: Bulk Upload
    ("06-bulk-upload.png",
     "Bulk Upload via Excel/CSV",
     ["Download template with format spec", "Upload hundreds of MCQs at once", "Row-level validation & error reporting", "Auto-mapping of tech stacks & topics"],
     "For large-scale content creation, SMEs can bulk upload MCQs via Excel or CSV. Download our template, fill in your questions, and upload. The system validates format, maps tech stacks automatically, and reports row-level errors for any issues."),

    # Scene 6: My Questions
    ("02-my-questions.png",
     "My Questions — Filter & Search",
     ["Filter by status, tech stack, topic, difficulty", "Keyword search across questions", "Sort by date, track lifecycle", "Edit drafts, view review feedback"],
     "The 'My Questions' view shows all MCQs you've created with full filtering — by status, tech stack, topic, and difficulty. Search by keywords, sort by date, and track the lifecycle of each question from draft to approval."),

    # Scene 7: Pending Reviews
    ("05-pending-reviews.png",
     "Review Workflow",
     ["Reviewer queue with pending MCQs", "AI Copilot quality analysis", "Approve/Reject with mandatory feedback", "Comment threads on questions"],
     "Admins and assigned reviewers see pending MCQs in their review queue. They can examine each question, run an AI quality analysis, and approve or reject with mandatory feedback. The AI Copilot assists reviewers by analyzing question quality, detecting potential issues, and providing confidence scores."),

    # Scene 8: Question Bank
    ("15-question-bank.png",
     "Question Bank — Admin View",
     ["All MCQs across all creators", "Advanced filters & export to Excel", "Bulk operations", "Complete audit trail"],
     "Admins have access to the full Question Bank — all MCQs across all creators. They can filter, search, export to Excel, manage users, and view the complete audit trail of every action taken on the platform."),

    # Scene 9: Analytics
    ("09-analytics.png",
     "Analytics Dashboard",
     ["MCQ distribution by tech stack", "Status breakdown charts", "Reviewer leaderboards", "Date-range filtering & export"],
     "The Analytics dashboard provides visual insights — MCQ distribution by tech stack, status breakdown charts, reviewer leaderboards, and approval rate metrics. Data can be filtered by date range and exported."),

    # Scene 10: Kanban
    ("09-kanban.png",
     "Kanban Board — Visual Workflow",
     ["Drag-and-drop across status columns", "Draft → Ready → Approved pipeline", "Visual progress tracking", "Quick status updates"],
     "The Kanban board gives a visual workflow view — drag and drop MCQs across status columns from Draft to Approved."),

    # Scene 11: Live Quiz
    ("live-01-host-lobby.png",
     "Live Quiz Sessions — Host View",
     ["Create real-time quiz sessions", "Generate shareable PIN code", "WebSocket for instant updates", "Timer sync across all devices"],
     "QuizHub supports real-time live quiz sessions. A host creates a session, selects questions, and generates a PIN code. Participants join using the PIN and answer questions in real-time with a live leaderboard."),

    # Live Quiz Player
    ("live-02-player-join.png",
     "Live Quiz — Player Join",
     ["Join with session PIN", "Real-time question delivery", "Instant answer feedback", "Live leaderboard updates"],
     "Players join using the PIN and see questions in real-time. The live session uses WebSocket connections — timer sync, answer submissions, and leaderboard changes happen instantly across all connected devices."),

    # Live Quiz Results
    ("live-05-final-results.png",
     "Live Quiz — Final Results",
     ["Final leaderboard with scores", "Question-by-question breakdown", "Performance analytics", "Session history saved"],
     "At the end of the session, the final leaderboard shows scores, rankings, and a question-by-question breakdown of each player's performance."),

    # Scene 12: Smart Interview Kit
    ("10-smart-interview-kit.png",
     "Smart Interview Kit — AI Resume Analysis",
     ["Upload candidate resume (PDF)", "AI generates 30 personalized questions", "6 categories: Technical, Coding, SQL, Project, Behavioral, Scenario", "Model answers included"],
     "The Smart Interview Kit is our AI-powered interview preparation tool. Upload a candidate's resume, and the AI generates 30 personalized interview questions across 6 categories — Technical, Coding, SQL, Project-Based, Behavioral, and Scenario questions — each with model answers. Questions are tailored to the candidate's actual skills and experience."),

    # Scene 13: Inbox
    ("12-inbox.png",
     "Inbox & Notifications",
     ["Real-time notification bell", "Review decisions & assignments", "Direct messaging between users", "Group chat for team collaboration"],
     "The platform includes real-time group chat for team collaboration, a notification system that alerts users about review decisions and assignments, and a full inbox for direct messaging between users."),

    # Scene 14: User Management
    ("16-user-management.png",
     "User Management — Admin Controls",
     ["Create/edit/deactivate users", "Assign roles (SME/Admin)", "Tech stack assignments per user", "Bulk user operations"],
     "Admins manage all platform users — creating accounts, assigning roles, mapping tech stacks to SMEs, and deactivating accounts when needed."),

    # Scene 15: Audit Log
    ("17-audit-log.png",
     "Audit Log — Complete Traceability",
     ["Every action logged with timestamp", "User, action type, entity tracking", "Filterable & searchable", "Compliance-ready audit trail"],
     "The Audit Log provides complete traceability — every action on the platform is logged with timestamp, user, action type, and affected entity. Fully filterable and searchable for compliance."),

    # Scene 16: Master Data
    ("18-master-data.png",
     "Master Data Management",
     ["Manage tech stacks & topics", "Add/edit/archive categories", "Hierarchical topic structure", "Used across all MCQ creation"],
     "Admins manage the master data — tech stacks, topics, and categories that form the foundation of all MCQ organization."),

    # Scene 17: i18n
    ("22-register.png",
     "Internationalization — 7 Languages",
     ["English, Hindi, French, German", "Kannada, Telugu, Urdu", "Full UI translation", "Accessible to global teams"],
     "QuizHub supports 7 languages — English, Hindi, French, German, Kannada, Telugu, and Urdu — making it accessible to global teams."),

    # Scene 18: Coding Questions
    ("21-coding-question.png",
     "Coding Questions",
     ["Code editor with syntax highlighting", "Multiple language support", "Test case validation", "AI-assisted generation"],
     "Beyond MCQs, QuizHub supports coding questions with a built-in code editor, syntax highlighting, and test case validation."),

    # Scene 19: Leaderboard
    ("23-leaderboard.png",
     "Leaderboard & Gamification",
     ["SME contribution rankings", "Points for creation & reviews", "Monthly/all-time views", "Encourages quality content creation"],
     "The leaderboard gamifies the MCQ creation process — SMEs earn points for creating and reviewing questions, encouraging quality content creation."),

    # Closing slide
    (None, None, None, None),  # handled specially
]


def add_title_slide(prs):
    """Add the opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Purple gradient background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x7B, 0x2D, 0x8B)  # Accenture purple

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "QuizHub AI"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Intelligent MCQ Management Platform"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(220, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Tagline
    p3 = tf.add_paragraph()
    p3.text = ""
    p4 = tf.add_paragraph()
    p4.text = "Accenture Hack-N-Stack 2026"
    p4.font.size = Pt(20)
    p4.font.color.rgb = RGBColor(200, 200, 255)
    p4.alignment = PP_ALIGN.CENTER

    # Tech stack info
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p5 = tf2.paragraphs[0]
    p5.text = "Spring Boot 3.2 • React 19 • Spring AI • MySQL 8 • WebSocket"
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(180, 180, 220)
    p5.alignment = PP_ALIGN.CENTER
    p6 = tf2.add_paragraph()
    p6.text = "2,029 Tests | 92.5% Backend Coverage | 80% Frontend Coverage"
    p6.font.size = Pt(14)
    p6.font.color.rgb = RGBColor(160, 160, 200)
    p6.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "Welcome to the QuizHub AI demo. This is an intelligent MCQ management platform "
        "built for Accenture's Hack-N-Stack 2026. The platform is built with Spring Boot 3.2, "
        "React 19, Spring AI with LLM integration, MySQL 8, and WebSocket for real-time features. "
        "It has 2,029 automated tests with 92.5% backend coverage and 80% frontend coverage."
    )


def add_closing_slide(prs):
    """Add the closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x7B, 0x2D, 0x8B)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = ""
    p3 = tf.add_paragraph()
    p3.text = "350+ Features • 25+ Modules • AI-Powered"
    p3.font.size = Pt(24)
    p3.font.color.rgb = RGBColor(220, 220, 255)
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "2,029 Tests • Zero Failures • Production Ready"
    p4.font.size = Pt(20)
    p4.font.color.rgb = RGBColor(200, 200, 255)
    p4.alignment = PP_ALIGN.CENTER

    # Stats box
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p5 = tf2.paragraphs[0]
    p5.text = "Backend: 1,072 tests • 92.5% coverage (JaCoCo)"
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(180, 180, 220)
    p5.alignment = PP_ALIGN.CENTER
    p6 = tf2.add_paragraph()
    p6.text = "Frontend: 957 tests • 80.37% coverage (Jest)"
    p6.font.size = Pt(16)
    p6.font.color.rgb = RGBColor(180, 180, 220)
    p6.alignment = PP_ALIGN.CENTER

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "Thank you for watching the QuizHub AI demo. Built with Spring Boot, React 19, "
        "Spring AI, and MySQL — delivering intelligent MCQ management at scale. "
        "2,029 automated tests, zero failures, production-ready."
    )


def add_content_slide(prs, image_file, title, bullets, notes):
    """Add a content slide with screenshot + text overlay."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    img_path = os.path.join(SHOTS_DIR, image_file)
    if os.path.exists(img_path):
        # Place screenshot on the left side (60% width)
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.2), width=Inches(7.8))

    # Title bar at top
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Bullet points on the right
    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(8.3), Inches(1.5), Inches(4.8), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(220, 220, 255)
            p.space_after = Pt(12)

    # Speaker notes (voiceover)
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def main():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, slide_data in enumerate(SLIDES):
        if i == 0:
            add_title_slide(prs)
        elif i == len(SLIDES) - 1:
            add_closing_slide(prs)
        else:
            img, title, bullets, notes = slide_data
            add_content_slide(prs, img, title, bullets, notes)

    prs.save(OUTPUT)
    print(f"✅ PowerPoint saved: {OUTPUT}")
    print(f"   {len(SLIDES)} slides with screenshots & speaker notes")
    print(f"\n💡 To record with voiceover in PowerPoint:")
    print(f"   1. Open the .pptx file")
    print(f"   2. Go to Slide Show → Record Slide Show")
    print(f"   3. Read the speaker notes aloud (shown at bottom)")
    print(f"   4. Export as video: File → Export → Create a Video")


if __name__ == "__main__":
    main()
