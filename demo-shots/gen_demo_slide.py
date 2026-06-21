"""
Generate a single demo slide with PowerPoint animations.
This demonstrates how python-pptx + raw XML animations work.

Animations in PPTX are stored as XML in the slide's timing tree.
python-pptx doesn't have a high-level animation API, so we inject
animation XML directly into the slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── Create Presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# ─── Background: Dark gradient ───────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.gradient()
fill.gradient_stops[0].color.rgb = RGBColor(0x0F, 0x0F, 0x2D)  # Deep navy
fill.gradient_stops[0].position = 0.0
fill.gradient_stops[1].color.rgb = RGBColor(0x1A, 0x1A, 0x4E)  # Purple-navy
fill.gradient_stops[1].position = 1.0

# ─── Shape 1: Title (will fly in from left) ──────────────────────────────────
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🚀 QuizHub AI"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# ─── Shape 2: Subtitle (will fade in) ────────────────────────────────────────
sub_box = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(9), Inches(1))
tf2 = sub_box.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "AI-Powered Question Management & Live Quiz Platform"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0xA7, 0x8B, 0xFA)  # Purple accent
p2.alignment = PP_ALIGN.CENTER

# ─── Shape 3: Stats card (will zoom in) ──────────────────────────────────────
card = slide.shapes.add_shape(
    1,  # MSO_SHAPE.ROUNDED_RECTANGLE
    Inches(3.5), Inches(4.5), Inches(6), Inches(2)
)
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x3F)
card.line.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
card.line.width = Pt(2)
card.shadow.inherit = False

# Add text to the card
tf3 = card.text_frame
tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
p3 = tf3.paragraphs[0]
p3.text = "421 Features  •  2,029 Tests  •  92.5% Coverage"
p3.font.size = Pt(22)
p3.font.bold = True
p3.font.color.rgb = RGBColor(0x34, 0xD3, 0x99)  # Green accent
p3.alignment = PP_ALIGN.CENTER

# ─── Shape 4: Team info (will appear) ────────────────────────────────────────
team_box = slide.shapes.add_textbox(Inches(3.5), Inches(6.8), Inches(6), Inches(0.5))
tf4 = team_box.text_frame
p4 = tf4.paragraphs[0]
p4.text = "Valkey Hack-N-Stack 2026"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
p4.alignment = PP_ALIGN.CENTER

# ─── ANIMATIONS ──────────────────────────────────────────────────────────────
# PowerPoint animations use the p:timing XML element.
# Each shape needs a unique ID reference (spTgt) matching its shape ID.

title_id = str(title_box.shape_id)
sub_id = str(sub_box.shape_id)
card_id = str(card.shape_id)
team_id = str(team_box.shape_id)

# Animation XML template with:
# 1. Title: Fly In from Left (on click)
# 2. Subtitle: Fade In (after 0.5s)
# 3. Stats Card: Zoom In (after 0.5s)
# 4. Team: Appear (after 0.3s)

timing_xml = f'''
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <!-- Main sequence (triggered on click) -->
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <!-- Click 1: Title flies in from left -->
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="2" presetClass="entr" presetSubtype="8" fill="hold" nodeType="clickEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="1" fill="hold">
                                        <p:stCondLst>
                                          <p:cond delay="0"/>
                                        </p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl>
                                        <p:spTgt spid="{title_id}"/>
                                      </p:tgtEl>
                                      <p:attrNameLst>
                                        <p:attrName>style.visibility</p:attrName>
                                      </p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:anim calcmode="lin" valueType="num">
                                    <p:cBhvr additive="base">
                                      <p:cTn id="7" dur="500" fill="hold"/>
                                      <p:tgtEl>
                                        <p:spTgt spid="{title_id}"/>
                                      </p:tgtEl>
                                      <p:attrNameLst>
                                        <p:attrName>ppt_x</p:attrName>
                                      </p:attrNameLst>
                                    </p:cBhvr>
                                    <p:tavLst>
                                      <p:tav tm="0"><p:val><p:strVal val="0-#ppt_w/2"/></p:val></p:tav>
                                      <p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav>
                                    </p:tavLst>
                                  </p:anim>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
                <!-- After title: Subtitle fades in -->
                <p:par>
                  <p:cTn id="8" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="9" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="500"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="10" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="afterEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="11" dur="1" fill="hold">
                                        <p:stCondLst>
                                          <p:cond delay="0"/>
                                        </p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl>
                                        <p:spTgt spid="{sub_id}"/>
                                      </p:tgtEl>
                                      <p:attrNameLst>
                                        <p:attrName>style.visibility</p:attrName>
                                      </p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:animEffect transition="in" filter="fade">
                                    <p:cBhvr>
                                      <p:cTn id="12" dur="700"/>
                                      <p:tgtEl>
                                        <p:spTgt spid="{sub_id}"/>
                                      </p:tgtEl>
                                    </p:cBhvr>
                                  </p:animEffect>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
                <!-- After subtitle: Stats card zooms in -->
                <p:par>
                  <p:cTn id="13" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="14" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="500"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="15" presetID="53" presetClass="entr" presetSubtype="32" fill="hold" nodeType="afterEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="16" dur="1" fill="hold">
                                        <p:stCondLst>
                                          <p:cond delay="0"/>
                                        </p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl>
                                        <p:spTgt spid="{card_id}"/>
                                      </p:tgtEl>
                                      <p:attrNameLst>
                                        <p:attrName>style.visibility</p:attrName>
                                      </p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:animEffect transition="in" filter="fade">
                                    <p:cBhvr>
                                      <p:cTn id="17" dur="600"/>
                                      <p:tgtEl>
                                        <p:spTgt spid="{card_id}"/>
                                      </p:tgtEl>
                                    </p:cBhvr>
                                  </p:animEffect>
                                  <p:animScale>
                                    <p:cBhvr>
                                      <p:cTn id="18" dur="600" fill="hold"/>
                                      <p:tgtEl>
                                        <p:spTgt spid="{card_id}"/>
                                      </p:tgtEl>
                                    </p:cBhvr>
                                    <p:from x="50000" y="50000"/>
                                    <p:to x="100000" y="100000"/>
                                  </p:animScale>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
                <!-- After card: Team info appears -->
                <p:par>
                  <p:cTn id="19" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="20" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="300"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="21" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="afterEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="22" dur="1" fill="hold">
                                        <p:stCondLst>
                                          <p:cond delay="0"/>
                                        </p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl>
                                        <p:spTgt spid="{team_id}"/>
                                      </p:tgtEl>
                                      <p:attrNameLst>
                                        <p:attrName>style.visibility</p:attrName>
                                      </p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:animEffect transition="in" filter="fade">
                                    <p:cBhvr>
                                      <p:cTn id="23" dur="500"/>
                                      <p:tgtEl>
                                        <p:spTgt spid="{team_id}"/>
                                      </p:tgtEl>
                                    </p:cBhvr>
                                  </p:animEffect>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    <p:bldP spid="{title_id}" grpId="0"/>
    <p:bldP spid="{sub_id}" grpId="0"/>
    <p:bldP spid="{card_id}" grpId="0"/>
    <p:bldP spid="{team_id}" grpId="0"/>
  </p:bldLst>
</p:timing>
'''

# Parse and inject timing XML into the slide
timing_element = etree.fromstring(timing_xml)
slide._element.append(timing_element)

# ─── Save ─────────────────────────────────────────────────────────────────────
output_path = '/Users/veera.konjeti/Desktop/hack-n-stack/demo-shots/DEMO_SINGLE_SLIDE.pptx'
prs.save(output_path)
print(f"✅ Demo slide saved: {output_path}")
print()
print("How it works when you open in PowerPoint:")
print("  1. Click → Title '🚀 QuizHub AI' flies in from the left")
print("  2. Auto → Subtitle fades in (0.5s delay)")
print("  3. Auto → Stats card zooms in from center (0.5s delay)")
print("  4. Auto → Team info fades in (0.3s delay)")
print()
print("Animation types used:")
print("  • Fly In (presetID=2) — element slides from off-screen")
print("  • Fade (presetID=10) — element fades from transparent to opaque")
print("  • Grow/Zoom (presetID=53 + animScale) — element scales from 50% to 100%")
print()
print("If this looks good, I'll generate the full animated PPT!")
